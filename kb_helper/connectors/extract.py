"""Turn file bytes into plain text for the model. Best-effort: unsupported types yield ''."""
from __future__ import annotations

import csv
import io
import json
import re
from html.parser import HTMLParser
from pathlib import PurePosixPath

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml", ".xml",
    ".log", ".ini", ".cfg", ".toml", ".py", ".js", ".ts", ".sh", ".ps1", ".sql", ".bat",
}
HTML_EXTENSIONS = {".html", ".htm", ".aspx", ".xhtml"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | HTML_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}


class _HTMLText(HTMLParser):
    _SKIP = {"script", "style", "noscript", "template"}
    _BLOCK = {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6", "section", "article", "table"}

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP:
            self._skip_depth += 1
        elif tag in self._BLOCK:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP and self._skip_depth:
            self._skip_depth -= 1
        elif tag in self._BLOCK:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip_depth:
            self._parts.append(data)

    def text(self) -> str:
        return normalise_whitespace("".join(self._parts))


def normalise_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def html_to_text(html: str) -> str:
    parser = _HTMLText()
    parser.feed(html)
    parser.close()
    return parser.text()


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for number, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(f"[page {number}]\n{page_text}")
    return "\n\n".join(pages)


def _docx(data: bytes) -> str:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        style = (paragraph.style.name if paragraph.style is not None else "") or ""
        text = paragraph.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "1"
            parts.append("#" * min(int(level), 6) + " " + text)
        else:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            if any(cell.strip() for cell in cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {number}]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n".join(parts)


def _csv(data: bytes, delimiter: str) -> str:
    rows = csv.reader(io.StringIO(_decode(data)), delimiter=delimiter)
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)


def extension_of(filename: str) -> str:
    return PurePosixPath(filename).suffix.lower()


def is_supported(filename: str) -> bool:
    return extension_of(filename) in SUPPORTED_EXTENSIONS


def extract_text(data: bytes, filename: str, content_type: str | None = None) -> str:
    """Extract readable text from ``data``. Returns '' when the format is not supported."""
    ext = extension_of(filename)
    ctype = (content_type or "").split(";")[0].strip().lower()
    try:
        if ext == ".pdf" or ctype == "application/pdf":
            return normalise_whitespace(_pdf(data))
        if ext == ".docx" or ctype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return normalise_whitespace(_docx(data))
        if ext == ".xlsx" or ctype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return normalise_whitespace(_xlsx(data))
        if ext == ".pptx" or ctype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return normalise_whitespace(_pptx(data))
        if ext in HTML_EXTENSIONS or ctype == "text/html":
            return html_to_text(_decode(data))
        if ext == ".csv":
            return _csv(data, ",")
        if ext == ".tsv":
            return _csv(data, "\t")
        if ext == ".json" or ctype == "application/json":
            try:
                return json.dumps(json.loads(_decode(data)), indent=2, ensure_ascii=False)
            except ValueError:
                return normalise_whitespace(_decode(data))
        if ext in TEXT_EXTENSIONS or ctype.startswith("text/"):
            return normalise_whitespace(_decode(data))
    except Exception as exc:  # noqa: BLE001 - extraction libraries raise many types
        return f"[could not extract text from {filename}: {exc.__class__.__name__}: {exc}]"
    return ""
